import os
import io
import re
import math
import docx
from pypdf import PdfReader
from pptx import Presentation
import mammoth
import pytesseract
from PIL import Image
try:
    from moviepy import VideoFileClip
except ImportError:
    from moviepy.editor import VideoFileClip
import tempfile

try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

# Anclas por categoría (RGB 0-255) ampliadas para fuentes / highlights / PDFs
# Incluye colores de text (font) típicos del panel de Word, más saturados/oscuros
_ANCHOR_RGB = {
    "VERB": [
        (0x92, 0xD0, 0x50), (0x00, 0xFF, 0x00), (0x00, 0xB0, 0x50),  # highlight
        (0x70, 0xAD, 0x47), (0x54, 0x82, 0x35), (0x00, 0x99, 0x00),  # font greens
        (0xD9, 0xEA, 0xD3), (0xB6, 0xD7, 0xA8), (0x93, 0xC4, 0x7D),  # google docs/pastel greens
    ],
    "NOUN": [
        (0x00, 0x70, 0xC0), (0x00, 0xB0, 0xF0), (0x4F, 0x81, 0xBD),  # highlight
        (0x44, 0x72, 0xC4), (0x2E, 0x75, 0xB6), (0x1F, 0x4E, 0x79),  # font blues
        (0x00, 0x00, 0xFF), (0x00, 0x00, 0xCC),                       # vivid blues
        (0xC9, 0xDA, 0xF8), (0x9F, 0xC5, 0xE8), (0xD0, 0xE0, 0xE3),  # pastels nouns/blues
        (0x6D, 0x9E, 0xEB), (0xA4, 0xC2, 0xF4), (0x6F, 0xA8, 0xDC), (0x3C, 0x78, 0xD8), # GD blues
        (0x4A, 0x86, 0xE8), # Cristofer blue
    ],
    "ADJ": [
        (0xFF, 0xFF, 0x00), (0xFF, 0xC0, 0x00), (0xFF, 0xD9, 0x66),  # highlight
        (0xED, 0x7D, 0x31), (0xFF, 0x99, 0x00), (0xFF, 0x80, 0x00),  # font oranges
        (0xFF, 0xF2, 0xCC), (0xFC, 0xE5, 0xCD),                       # pastels adjectives
    ],
    "ADV": [
        (0xFF, 0x00, 0xFF), (0xE2, 0x6B, 0x0A), (0xFF, 0x00, 0x66),  # highlight
        (0xFF, 0x00, 0x00), (0xC0, 0x00, 0x00), (0xFF, 0x14, 0x93),  # reds/pinks
        (0xDC, 0x14, 0x3C), (0xFF, 0x69, 0xB4), (0xFC, 0x0F, 0xC0),  # fuchsia/hot pink
        (0xEA, 0xD1, 0xDC), (0xF4, 0xCC, 0xCC), (0xFF, 0x66, 0xFF),  # pastels adverbs/pinks
    ],
    "DET": [
        (0xA6, 0xA6, 0xA6), (0x7F, 0x7F, 0x7F), (0xBF, 0xBF, 0xBF),  # highlight
        (0x59, 0x59, 0x59), (0x40, 0x40, 0x40), (0x80, 0x80, 0x80),  # font grays
        (0xB7, 0xB7, 0xB7), (0x99, 0x99, 0x99), (0xD9, 0xD9, 0xD9),  # pastels grays
        (0xCC, 0xCC, 0xCC), (0xEF, 0xEF, 0xEF),                       # GD Grays
        (0xF2, 0xF2, 0xF2), (0x26, 0x26, 0x26), (0x66, 0x66, 0x66), (0xAE, 0xAA, 0xAA), (0xD6, 0xDC, 0xE4),
        (0xBD, 0xDB, 0xDB), (0x9E, 0x9E, 0x9E), (0xB0, 0xB0, 0xB0),
        (0xBD, 0xDB, 0xDB), (0x9E, 0x9E, 0x9E), (0xB0, 0xB0, 0xB0), (0xCC, 0xCC, 0xCC), (0xD9, 0xD9, 0xD9),
    ],
    "ADP": [
        (0x70, 0x30, 0xA0), (0x60, 0x49, 0x7A), (0xB1, 0xA0, 0xC7),  # highlight
        (0x80, 0x64, 0xA2), (0x5C, 0x2D, 0x91), (0x9C, 0x27, 0xB0),  # font purples
        (0x84, 0x3F, 0xA0), (0x68, 0x21, 0x7A), (0x6A, 0x1B, 0x9A),  # standard dark purples
        (0xB4, 0xA7, 0xD6), (0xD9, 0xD2, 0xE9), (0x99, 0x66, 0xFF),  # pastels prepositions
        (0x99, 0x00, 0xFF), # Cristofer purple
        (0xE4, 0xDF, 0xEC), (0xCC, 0xC1, 0xDA), (0x4F, 0x34, 0x62), (0xCC, 0xC0, 0xDA), (0x4A, 0x14, 0x8C),
        (0xAB, 0x47, 0xBC), (0x8E, 0x24, 0xAA), (0x7B, 0x1F, 0xA2),
    ],
}

def _color_distance(r1, g1, b1, r2, g2, b2):
    """Distancia euclidiana en RGB normalizada (0-1)."""
    return math.sqrt(((r1 - r2) / 255) ** 2 + ((g1 - g2) / 255) ** 2 + ((b1 - b2) / 255) ** 2)

def rgb_to_category(r, g, b, max_distance=0.35):
    """
    Mapea RGB a categoría gramatical por proximidad a anclas.
    Retorna (category, confidence) o (None, 0). confidence en [0, 1].
    """
    if r is None or g is None or b is None:
        return None, 0.0
    r, g, b = int(r), int(g), int(b)
    best_cat, best_conf = None, 0.0
    for cat, anchors in _ANCHOR_RGB.items():
        for (ar, ag, ab) in anchors:
            d = _color_distance(r, g, b, ar, ag, ab)
            conf = max(0, 1.0 - (d / max_distance))
            if conf > best_conf:
                best_conf = conf
                best_cat = cat
    return (best_cat, round(best_conf, 2)) if best_conf >= 0.5 else (None, 0.0)

def hex_to_category(hex_color):
    """
    Maps hex colors from Word to grammatical categories (exact match).
    Includes both highlight/shading palette AND the standard Word font-color palette
    (Office theme colors and standard colors panel).
    """
    if not hex_color: return None
    hex_color = str(hex_color).strip().upper().replace("#", "")
    # ---- VERB: Verbo (verde) ----
    # Highlight/shading greens + Office theme greens (font color panel)
    if hex_color in [
        '92D050', '00FF00', '76933C', 'C4D79B', 'EBF1DE', '00B050',  # highlight/shading
        '70AD47', '548235', '375623', 'A9D18E', 'E2EFDA',             # Office theme greens
        '00B04F', '009900', '33CC00', '00CC00', '66FF00',             # other greens
        'D9EAD3', 'B6D7A8', '93C47D',                                 # Google Docs greens
    ]: return "VERB"
    # ---- NOUN: Sustantivo (azul) ----
    if hex_color in [
        '0070C0', '00B0F0', '4F81BD', 'B8CCE4', 'DBE5F1', 'B7DEE8', '31859B', '92CDDC',  # highlight
        '4472C4', '2E75B6', '1F4E79', 'BDD7EE', 'DEEAF1', '2F5496',  # Office theme blues
        '0000FF', '003DFF', '0563C1', '17375E', '003399',             # other blues
        '00AEEF', '0096C8', '4BACC6',                                 # teal/sky blues
        'C9DAF8', '9FC5E8', 'D0E0E3', '6D9EEB', 'A4C2F4', '6FA8DC', '3C78D8', '4A86E8', # GD blues
    ]: return "NOUN"
    # ---- ADJ: Adjetivo (amarillo/naranja) ----
    if hex_color in [
        'FFFF00', 'FFC000', 'FFD966', 'FFF2CC', 'FCD5B4',  # highlight
        'ED7D31', 'C55A11', '843C0C', 'F4B183', 'FCE4D6',  # Office orange theme
        'FFF100', 'FFE135', 'FFD700', 'F9A825', 'F57F17',  # other yellows
        'FF9900', 'FF8000', 'F7A239', 'E6A118',            # ambers/oranges
        'FCE5CD',                                           # Google Docs orange light
    ]: return "ADJ"
    # ---- ADV: Adverbio (rosa/magenta/fucsia) ----
    if hex_color in [
        'FF00FF', 'FF0066', 'E26B0A', 'F2DCDB', 'FDEADA', 'E46C0A',  # highlight
        'FF0000', 'C00000', 'D00000', 'A50021', '9B2335',             # reds
        'FF1493', 'FF69B4', 'FF66CC', 'F50057', 'E91E63',            # pinks/fuchsia
        'EA3560', 'FC0FC0', 'DC143C', 'FF4081', 'FF80AB',            # hot pinks
        'C0004B', 'EAD1DC', 'F4CCCC', 'FF66FF',                       # dark magenta + GD pinks
    ]: return "ADV"
    # ---- DET: Determinante / Artículo / Pronombre (gris) ----
    if hex_color in [
        'A6A6A6', 'BFBFBF', 'D9D9D9', 'F2F2F2', '7F7F7F',  # highlight
        '595959', '262626', '404040', '666666', '808080',    # Office dark grays
        'AEAAAA', 'D6DCE4', 'BDBDBD', '9E9E9E', 'B0B0B0',  # light-medium grays
        'B7B7B7', '999999', 'CCCCCC', 'EFEFEF',             # custom grays
    ]: return "DET"
    # ---- ADP: Preposición (morado/violeta) ----
    if hex_color in [
        '7030A0', 'B1A0C7', 'E4DFEC', 'CCC1DA', '60497A',  # highlight
        '8064A2', '4F3462', 'CCC0DA', 'E4DFEC', '843FA0',  # Office purple theme
        '5C2D91', '68217A', '6A1B9A', '4A148C', '9C27B0',  # deep purples
        'AB47BC', '8E24AA', '7B1FA2',                       # medium purples
        'B4A7D6', 'D9D2E9', '9966FF', '9900FF',             # custom/GD purples
    ]: return "ADP"
    return None

def color_to_category(hex_or_rgb, default_confidence=0.85):
    """
    Convierte color (hex string o tuple (r,g,b)) a (category, confidence).
    Primero intenta hex exacto; si no, RGB fuzzy.
    """
    if hex_or_rgb is None:
        return None, 0.0
    if isinstance(hex_or_rgb, (list, tuple)) and len(hex_or_rgb) >= 3:
        r, g, b = hex_or_rgb[0], hex_or_rgb[1], hex_or_rgb[2]
        cat = hex_to_category(f"{int(r):02X}{int(g):02X}{int(b):02X}")
        if cat:
            return cat, default_confidence
        return rgb_to_category(r, g, b)
    h = str(hex_or_rgb).strip().replace("#", "")
    if len(h) >= 6:
        cat = hex_to_category(h)
        if cat:
            return cat, 1.0
        try:
            r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
            return rgb_to_category(r, g, b)
        except Exception:
            pass
    return None, 0.0

def normalize_highlight(word, expected_category, source="format_color", confidence=1.0, color_hex=None):
    """Formato unificado de highlight para cualquier extractor."""
    return {
        "word": word,
        "expected_category": expected_category,
        "source": source,
        "confidence": min(1.0, max(0.0, float(confidence))),
        "color_hex": color_hex,
    }

def get_run_color_category(run):
    """
    Detects grammatical category from a Word run using three methods (in priority order):
    1. Highlight color index (wdColorIndex enum).
    2. Shading fill color (paragraph/run background).
    3. Font color (run.font.color.rgb) — used when no highlight/shading is set.
    Returns (category, source) or (None, None).
    """
    highlights = {
        4: 'VERB', 11: 'VERB',         # Verdes (Brillante y Estándar)
        2: 'NOUN', 3: 'NOUN', 9: 'NOUN', 10: 'NOUN', # Azules (Azul, Turquesa, Azul oscuro, Teal)
        7: 'ADJ', 14: 'ADJ', 1: 'ADJ',  # Amarillos (Amarillo, Amarillo oscuro y Legacy 1)
        5: 'ADV', 6: 'ADV',            # Rosas y Rojos (Pink, Red)
        12: 'ADP', 13: 'ADP',          # Morados (Violeta, Rojo oscuro/Bordó)
        16: 'DET', 15: 'DET'           # Grises
    }
    # 1. Highlight color
    if run.font.highlight_color in highlights:
        return highlights[run.font.highlight_color]
    # 2. Shading / background fill
    try:
        rPr = run._element.get_or_add_rPr()
        shd = rPr.xpath('w:shd')
        if shd:
            fill = shd[0].get('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}fill')
            if fill and fill != 'auto':
                cat = hex_to_category(fill)
                if cat:
                    return cat
    except: pass
    # 3. Font color (run.font.color.rgb)
    try:
        from docx.shared import RGBColor
        fc = run.font.color
        if fc and fc.type is not None:
            rgb = fc.rgb  # RGBColor instance
            if rgb is not None:
                r, g, b = rgb[0], rgb[1], rgb[2]
                if not _is_neutral_color(r, g, b):
                    cat = hex_to_category(f"{r:02X}{g:02X}{b:02X}")
                    if cat:
                        return cat
                    cat, conf = rgb_to_category(r, g, b)
                    if cat:
                        return cat
    except: pass
    return None

# Título esperado según actividad M2 AI2
TITULO_RELATO_ESPERADO = "el relato de mi historia"

# Colores para vista previa (verbos verde, sustantivos azul, adjetivos amarillo, adverbios rosa, artículos gris, preposiciones morado)
CATEGORY_CSS = {
    "VERB": "#92D050",
    "NOUN": "#0070C0",
    "ADJ": "#FFC000",
    "ADV": "#FF00FF",
    "DET": "#A6A6A6",
    "ADP": "#7030A0",
}

def extract_audio_url_from_text(text):
    """Extrae la primera URL de audio/video del texto (típicamente al final del documento)."""
    if not text or not isinstance(text, str):
        return None
    url_pattern = re.compile(r'https?://[^\s<>"\']+', re.IGNORECASE)
    matches = url_pattern.findall(text.strip())
    for url in reversed(matches):
        u = url.rstrip('.,;:)')
        if any(x in u.lower() for x in ('drive', 'dropbox', 'onedrive', '1drv', 'sharepoint', 'youtu', 'soundcloud', 'audio', 'video', 'mp3', 'mp4', 'watch', 'sharing')):
            return u
    return matches[-1] if matches else None

def detect_titulo_presente(text):
    """Indica si el título 'El relato de mi historia' está presente en el texto."""
    if not text:
        return False
    normalized = re.sub(r'\s+', ' ', text.strip().lower())
    return TITULO_RELATO_ESPERADO in normalized

def _check_fuente_arial_12_docx(doc):
    """Comprueba si el documento usa Arial 12pt (muestreo en primeros párrafos)."""
    try:
        count_arial, count_total = 0, 0
        for para in list(doc.paragraphs)[:30]:
            for run in para.runs:
                if not run.text.strip():
                    continue
                count_total += 1
                size_ok = (run.font.size is not None and run.font.size.pt >= 11.5 and run.font.size.pt <= 12.5)
                name_ok = (run.font.name and 'arial' in (run.font.name or '').lower())
                if size_ok or name_ok:
                    count_arial += 1
        if count_total == 0:
            return True
        return (count_arial / count_total) >= 0.5
    except Exception:
        return True

def _get_run_color_info(run):
    """
    Returns (category, source_label, color_hex) from a run, checking in order:
    1. Highlight color index.
    2. Shading fill (background).
    3. Font color — reads w:color w:val from XML directly.
       This works for BOTH RGB colors AND theme colors, because Word always
       writes the resolved hex value into the w:val attribute.
    """
    highlights = {
        4: 'VERB', 11: 'VERB',         # Verdes (Brillante y Estándar)
        2: 'NOUN', 3: 'NOUN', 9: 'NOUN', 10: 'NOUN', # Azules (Azul, Turquesa, Azul oscuro, Teal)
        7: 'ADJ', 14: 'ADJ', 1: 'ADJ',  # Amarillos (Amarillo, Amarillo oscuro y Legacy 1)
        5: 'ADV', 6: 'ADV',            # Rosas y Rojos (Pink, Red)
        12: 'ADP', 13: 'ADP',          # Morados (Violeta, Rojo oscuro/Bordó)
        16: 'DET', 15: 'DET'           # Grises
    }
    if run.font.highlight_color in highlights:
        return highlights[run.font.highlight_color], "docx_highlight", None
    try:
        rPr = run._element.get_or_add_rPr()
        shd = rPr.xpath('w:shd')
        if shd:
            fill = shd[0].get('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}fill')
            if fill and fill != 'auto':
                cat = hex_to_category(fill)
                if cat:
                    return cat, "docx_shading", fill
    except: pass
    # Font color: read w:val from w:color XML element.
    # This approach works for ALL Word color types (explicit RGB, theme color, or tinted).
    # python-docx's fc.rgb raises TypeError for theme colors, so we bypass it by reading
    # the XML attribute directly — Word always resolves and writes the display hex into w:val.
    try:
        from docx.oxml.ns import qn
        rPr = run._element.get_or_add_rPr()
        color_elem = rPr.find(qn('w:color'))
        if color_elem is not None:
            val = color_elem.get(qn('w:val'))
            if val and val.upper() not in ('AUTO', '000000', 'FFFFFF', ''):
                val_upper = val.upper().replace('#', '')
                if len(val_upper) == 6:
                    cat = hex_to_category(val_upper)
                    if not cat:
                        r, g, b = int(val_upper[0:2], 16), int(val_upper[2:4], 16), int(val_upper[4:6], 16)
                        if not _is_neutral_color(r, g, b):
                            cat, _ = rgb_to_category(r, g, b, max_distance=0.42)
                    if cat:
                        return cat, "docx_font_color", val_upper
    except: pass
    return None, None, None


def extract_text_from_docx(file_bytes):
    """Extracts text, HTML, and highlights from a Word document."""
    doc = docx.Document(io.BytesIO(file_bytes))
    full_text = []
    highlights = []
    
    for para in doc.paragraphs:
        full_text.append(para.text)
        for run in para.runs:
            category, source, color_hex = _get_run_color_info(run)
            if category and run.text.strip():
                words = re.findall(r'\b[\wáéíóúÁÉÍÓÚñÑ]+\b', run.text)
                for w in words:
                    highlights.append(normalize_highlight(w, category, source=source, confidence=1.0, color_hex=color_hex))
    
    result = mammoth.convert_to_html(io.BytesIO(file_bytes))
    html = result.value
    full_str = "\n".join(full_text)
    titulo_ok = detect_titulo_presente(full_str)
    fuente_ok = _check_fuente_arial_12_docx(doc)
    audio_url = extract_audio_url_from_text(full_str)
    meta = {"titulo_ok": titulo_ok, "fuente_ok": fuente_ok, "audio_url": audio_url}
    html_with_highlights = _build_html_with_highlights(doc)
    meta["html_with_highlights"] = html_with_highlights if html_with_highlights else html
    return full_str, html, highlights, meta


def _build_html_with_highlights(doc):
    """Genera HTML que conserva los colores de resalte del Word (clases de palabras)."""
    parts = ['<div style="font-family: Arial, sans-serif; font-size: 12pt; padding: 1em;">']
    for para in doc.paragraphs:
        if not para.text.strip():
            parts.append("<p><br></p>")
            continue
        parts.append("<p>")
        for run in para.runs:
            if not run.text:
                continue
            category, source, _ = _get_run_color_info(run)
            escaped = run.text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            if category and category in CATEGORY_CSS:
                color = CATEGORY_CSS[category]
                if source == "docx_font_color":
                    # Font color: show with colored underline and text color (no background)
                    parts.append(f'<span style="color:{color}; border-bottom: 2px solid {color}; font-weight:600;">{escaped}</span>')
                else:
                    # Highlight or shading: show with background
                    parts.append(f'<span style="background-color:{color}; padding: 0 1px;">{escaped}</span>')
            else:
                parts.append(escaped)
        parts.append("</p>")
    parts.append("</div>")
    return "".join(parts)


def detect_sections(full_text):
    """
    Detecta automáticamente relato, reflexión y enlace de audio.
    Retorna dict con relato, reflexion, audio_url.
    """
    if not full_text or not isinstance(full_text, str):
        return {"relato": "", "reflexion": "", "audio_url": None}
    text = full_text.strip()
    audio_url = extract_audio_url_from_text(text)
    url_pos = text.find(audio_url) if audio_url else -1
    # Posibles inicios del párrafo de reflexión (últimas 800 caracteres suelen contener reflexión + enlace)
    tail = text[-1200:] if len(text) > 1200 else text
    reflexion_start_markers = [
        "diferencias que encontré",
        "diferencias entre",
        "expresarse de manera oral",
        "oral y escrita",
        "comunicación oral y escrita",
        "reflexión",
        "en mi caso",
        "en conclusión",
    ]
    reflexion_start = -1
    for m in reflexion_start_markers:
        idx = tail.lower().find(m.lower())
        if idx >= 0:
            reflexion_start = len(text) - len(tail) + idx
            break
    if reflexion_start >= 0:
        relato = text[:reflexion_start].strip()
        reflexion = text[reflexion_start:].strip()
    elif url_pos >= 0:
        relato = text[:url_pos].strip()
        reflexion = text[url_pos:].strip()
    else:
        relato = text
        reflexion = ""
    return {"relato": relato, "reflexion": reflexion, "audio_url": audio_url}

def _pdf_span_color_to_rgb(span_color):
    """Convierte color de PyMuPDF (int o 3 floats 0-1) a (r,g,b)."""
    if span_color is None:
        return None
    if isinstance(span_color, (list, tuple)) and len(span_color) >= 3:
        c = span_color[:3]
        if all(isinstance(x, (int, float)) for x in c):
            if all(0 <= x <= 1 for x in c):
                return (int(c[0] * 255), int(c[1] * 255), int(c[2] * 255))
            return (int(c[0]), int(c[1]), int(c[2]))
    if isinstance(span_color, int):
        return ((span_color >> 16) & 0xFF, (span_color >> 8) & 0xFF, span_color & 0xFF)
    return None

def _is_neutral_color(r, g, b):
    """True for white, near-white, black or near-black — not a meaningful highlight color."""
    brightness = r * 0.299 + g * 0.587 + b * 0.114
    if brightness > 230:
        return True
    if brightness < 25:
        return True
    if abs(r - g) < 15 and abs(g - b) < 15 and abs(r - b) < 15:
        return brightness > 180 or brightness < 50
    return False

def _extract_highlights_from_pdf_drawings(file_bytes):
    """
    Ruta A: PyMuPDF — detecta rectángulos de color (sombreado/fondo) debajo del texto.
    Ignora completamente el color del texto (foreground).
    """
    if fitz is None:
        return []
    highlights = []
    try:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        for page in doc:
            colored_rects = []
            for drawing in page.get_drawings():
                fill = drawing.get("fill")
                if fill is None:
                    continue
                rgb = _pdf_span_color_to_rgb(fill)
                if rgb is None or _is_neutral_color(*rgb):
                    continue
                category, conf = color_to_category(rgb, default_confidence=0.85)
                if not category:
                    continue
                rect = drawing.get("rect")
                if rect is None:
                    continue
                colored_rects.append((fitz.Rect(rect), category, conf, rgb))

            if not colored_rects:
                continue

            words_on_page = page.get_text("words")
            for w in words_on_page:
                word_rect = fitz.Rect(w[:4])
                word_text = w[4].strip()
                if not word_text or not re.match(r'^[\wáéíóúÁÉÍÓÚñÑ]+$', word_text):
                    continue
                for (crect, cat, conf, rgb) in colored_rects:
                    if crect.intersects(word_rect):
                        overlap = crect & word_rect
                        if overlap.width > word_rect.width * 0.5 and overlap.height > word_rect.height * 0.3:
                            highlights.append(normalize_highlight(
                                word_text, cat, source="pdf_drawing_fill", confidence=conf,
                                color_hex=f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"
                            ))
                            break
        doc.close()
    except Exception:
        pass
    return highlights

def _extract_highlights_from_pdf_annotations(file_bytes):
    """Ruta B: anotaciones tipo highlight/markup en PDF (pypdf/PyMuPDF)."""
    highlights = []
    try:
        if fitz is None:
            return highlights
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        for page in doc:
            for annot in page.annots() or []:
                if annot.type[0] != 8:  # 8 = PDF_ANNOT_HIGHLIGHT
                    continue
                r = annot.rect
                words = page.get_text("words", clip=r)
                col = annot.colors.get("stroke") or annot.colors.get("fill")
                rgb = _pdf_span_color_to_rgb(col)
                category, conf = (None, 0.0)
                if rgb:
                    category, conf = color_to_category(rgb, default_confidence=0.75)
                if not category:
                    category = "NOUN"
                    conf = 0.5
                for w in words:
                    word = w[4].strip()
                    if re.match(r'^[\wáéíóúÁÉÍÓÚñÑ]+$', word):
                        highlights.append(normalize_highlight(
                            word, category, source="pdf_annotation", confidence=conf, color_hex=None
                        ))
        doc.close()
    except Exception:
        pass
    return highlights


def _extract_highlights_from_pdf_spans(file_bytes):
    """
    Ruta C: Lee el COLOR DE FUENTE de cada span de texto usando PyMuPDF rawdict.
    Esta es la técnica más efectiva para PDFs exportados desde Word, donde el texto
    coloreado se preserva como color del span (no como rectángulo de fondo).
    
    page.get_text("rawdict") retorna la estructura:
      blocks → lines → spans → { "text": …, "color": int_rgb, "bbox": (x0,y0,x1,y1) }
    
    El entero 'color' es 0xRRGGBB codificado como valor decimal.
    """
    if fitz is None:
        return []
    highlights = []
    try:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        for page in doc:
            blocks = page.get_text("rawdict", flags=fitz.TEXT_PRESERVE_WHITESPACE).get("blocks", [])
            for block in blocks:
                if block.get("type") != 0:  # 0 = text block
                    continue
                for line in block.get("lines", []):
                    for span in line.get("spans", []):
                        color_int = span.get("color", 0)
                        span_text = span.get("text", "").strip()
                        if not span_text:
                            continue
                        # Convertir entero RGB a componentes
                        rgb = _pdf_span_color_to_rgb(color_int)
                        if rgb is None or _is_neutral_color(*rgb):
                            continue
                        r, g, b = rgb
                        hex_str = f"{r:02X}{g:02X}{b:02X}"
                        # Primero intenta mapeo exacto, luego fuzzy con umbral amplio
                        cat = hex_to_category(hex_str)
                        if not cat:
                            cat, _ = rgb_to_category(r, g, b, max_distance=0.42)
                        if not cat:
                            continue
                        # Extraer palabras del span y registrarlas
                        words = re.findall(r'\b[\wáéíóúÁÉÍÓÚñÑ]+\b', span_text)
                        for w in words:
                            highlights.append(normalize_highlight(
                                w, cat, source="pdf_font_color", confidence=0.85,
                                color_hex=hex_str
                            ))
        doc.close()
    except Exception:
        pass
    return highlights


def _build_html_from_pdf(file_bytes):
    """
    Genera HTML con vista previa coloreada de un PDF de forma GRANULAR (palabra por palabra).
    Resuelve el problema de PDFs con bloques de texto comprimidos donde los resaltados
    se perdían en bboxes de span demasiado grandes.
    """
    if fitz is None:
        return None
    try:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        parts = ['<div style="font-family: Arial, sans-serif; font-size: 11.5pt; padding: 1.5em; line-height: 1.6; color: #333;">']
        
        for page_num, page in enumerate(doc):
            if page_num > 0:
                parts.append('<hr style="border: 0; border-top: 1px dashed #ccc; margin: 2em 0;">')
            
            # 1. Obtener rectángulos de dibujo (fondos coloreados)
            colored_rects = []
            for drawing in page.get_drawings():
                fill = drawing.get("fill")
                if fill is None: continue
                rgb = _pdf_span_color_to_rgb(fill)
                if rgb and not _is_neutral_color(*rgb):
                    rect = drawing.get("rect")
                    if rect:
                        cat_bg, _ = color_to_category(rgb)
                        if cat_bg:
                            colored_rects.append((fitz.Rect(rect), cat_bg))

            # 2. Obtener TODAS las palabras con sus bboxes individuales
            # Esto es clave para la granularidad.
            words_info = page.get_text("words") # (x0, y0, x1, y1, "word", block_no, line_no, word_no)
            
            # 3. Categorizar cada palabra (por fondo o por fuente si fitz diera color por palabra)
            # Como get_text("words") no da color, lo sacaremos del dict en el paso siguiente
            # Pero usaremos las bboxes de las palabras para el cruce con los fondos.
            
            # 4. Obtener estructura de bloques para reconstruir el layout
            d = page.get_text("dict")
            blocks = d.get("blocks", [])
            
            if not blocks:
                text_plain = page.get_text("text")
                if text_plain: parts.append(f"<p>{text_plain}</p>")
                continue

            for b in blocks:
                if b.get("type") != 0: continue
                parts.append("<p style='margin-bottom: 0.8em;'>")
                
                for line in b.get("lines", []):
                    # En lugar de solo iterar spans, vamos a procesar palabras que caen en este span
                    for span in line.get("spans", []):
                        span_text = span.get("text", "")
                        if not span_text.strip():
                            parts.append(span_text) # Espacios, etc.
                            continue
                        
                        # Extraer palabras del span_text (respetando separadores si es posible)
                        # Pero lo más preciso es buscar qué palabras de 'words_info' intersectan este span
                        span_rect = fitz.Rect(span.get("bbox"))
                        
                        # Buscamos palabras que pertenezcan a este bloque/línea o intersecten el rect
                        span_words = [w for w in words_info if span_rect.intersects(fitz.Rect(w[:4]))]
                        
                        if not span_words:
                            # Fallback si no hay palabras detectadas por fitz (raro)
                            parts.append(span_text.replace("<", "&lt;").replace(">", "&gt;"))
                            continue
                        
                        # Reconstruir el span palabra por palabra con sus respectivos colores
                        # Nota: span_words puede tener más palabras de las contenidas en ESE span si el PDF es caótico
                        # Así que simplemente procesamos el span_text pero con lógica de subrayado inteligente
                        
                        # Mejor enfoque: El span_text ya tiene las palabras.
                        # Lo dividimos por espacios y a cada fragmento le buscamos su color.
                        sub_fragments = span_text.split(" ")
                        for i, frag in enumerate(sub_fragments):
                            clean_frag = frag.strip(",.;:()\"' ")
                            display_frag = frag.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                            
                            # Buscar color de FONDO para esta palabra específica
                            color_cat = None
                            if clean_frag:
                                # Buscar en words_info la palabra que coincida con la posición
                                # Usamos un pequeño margen para el matching espacial
                                for winfo in span_words:
                                    if winfo[4] == clean_frag:
                                        w_rect = fitz.Rect(winfo[:4])
                                        for crect, ccat in colored_rects:
                                            if crect.intersects(w_rect):
                                                overlap = crect & w_rect
                                                if overlap.width > w_rect.width * 0.4:
                                                    color_cat = ccat
                                                    break
                                        if color_cat: break
                                
                                # Si no hay fondo, buscar color de FUENTE del span completo
                                if not color_cat:
                                    color_int = span.get("color", 0)
                                    rgb_font = _pdf_span_color_to_rgb(color_int)
                                    if rgb_font and not _is_neutral_color(*rgb_font):
                                        r, g, b = rgb_font
                                        hex_str = f"{r:02X}{g:02X}{b:02X}"
                                        cat_f = hex_to_category(hex_str)
                                        if not cat_f:
                                            cat_f, _ = rgb_to_category(r, g, b, max_distance=0.42)
                                        if cat_f: color_cat = cat_f

                            if color_cat and color_cat in CATEGORY_CSS:
                                css_color = CATEGORY_CSS[color_cat]
                                # Usar background si es muy relevante o font-color + underline
                                # Para PDF, el background-color suele ser más fiel al "resaltado" (highlighter)
                                parts.append(f'<span style="background-color:{css_color}; padding: 0 1px; border-radius: 1px; font-weight:500;">{display_frag}</span>')
                            else:
                                parts.append(display_frag)
                            
                            if i < len(sub_fragments) - 1:
                                parts.append(" ")
                    parts.append(" ") # Espacio entre spans
                parts.append("</p>")
        
        parts.append("</div>")
        doc.close()
        return "".join(parts)
    except Exception:
        return None


def extract_text_from_pdf(file_bytes):
    """
    Extrae texto e highlights de un PDF usando PyMuPDF (fitz) únicamente.
    
    Devuelve:
      text: texto plano completo.
      html_preview: HTML coloreado para el visor de la app.
      highlights: lista de dicts con palabras detectadas por color.
    """
    if fitz is None:
        # Fallback ultra-seguro si fitz no está (no debería pasar en Docker)
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(file_bytes))
        text = "\n".join(p.extract_text() or "" for p in reader.pages)
        return text, f"<pre>{text}</pre>", []

    text = ""
    highlights = []
    html_preview = None

    try:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        
        # 1. Extraer texto plano (usando fitz para consistencia)
        for page in doc:
            text += page.get_text("text") + "\n"
        
        # 2. Rutas de Higgins (A, B, C)
        # Ruta A: Fondos/Drawings
        hl_a = _extract_highlights_from_pdf_drawings(file_bytes)
        # Ruta B: Anotaciones nativas
        hl_b = _extract_highlights_from_pdf_annotations(file_bytes)
        # Ruta C: Color de fuente (spans)
        hl_c = _extract_highlights_from_pdf_spans(file_bytes)

        # Acumular y deduplicar
        seen = set()
        for hl in hl_a + hl_b + hl_c:
            key = (hl["word"].lower(), hl["expected_category"])
            if key not in seen:
                seen.add(key)
                highlights.append(hl)
        
        doc.close()
    except Exception:
        pass

    # 3. Generar Vista Previa HTML
    try:
        html_preview = _build_html_from_pdf(file_bytes)
    except Exception:
        pass

    # Fallback final de seguridad si el HTML falló
    if not html_preview or len(html_preview) < 150:
        html_preview = f'<div style="font-family:sans-serif; padding:1em;"><pre style="white-space:pre-wrap;">{text}</pre></div>'

    return text, html_preview, highlights



def _pptx_run_highlight_rgb(run):
    """
    Obtiene (r,g,b) del highlight/shading de un run en pptx (color de FONDO, no de texto).
    Busca en el XML del run el elemento <a:highlight> o <a:solidFill> dentro de <a:rPr>.
    """
    try:
        rPr = run._r.get_or_add_rPr()
        ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
        # <a:highlight><a:srgbClr val="FFFF00"/></a:highlight>
        hl = rPr.findall(f"{{{ns}}}highlight")
        for h in hl:
            for srgb in h.findall(f"{{{ns}}}srgbClr"):
                val = srgb.get("val")
                if val and len(val) >= 6:
                    return (int(val[0:2], 16), int(val[2:4], 16), int(val[4:6], 16))
        # Fallback: shape fill behind text is not per-run in pptx
    except Exception:
        pass
    return None

def extract_text_from_pptx(file_bytes):
    """Extracts text, HTML, and highlights from PowerPoint (highlight/shading color, not font color)."""
    prs = Presentation(io.BytesIO(file_bytes))
    text_parts = []
    highlights = []
    html_parts = ['<div style="font-family: sans-serif; padding: 1em;">']
    for slide in prs.slides:
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
            for para in shape.text_frame.paragraphs:
                line = ""
                for run in para.runs:
                    t = run.text or ""
                    line += t
                    rgb = _pptx_run_highlight_rgb(run)
                    if rgb and t.strip() and not _is_neutral_color(*rgb):
                        category, conf = color_to_category(rgb, default_confidence=0.9)
                        if category:
                            words = re.findall(r'\b[\wáéíóúÁÉÍÓÚñÑ]+\b', t)
                            for w in words:
                                highlights.append(normalize_highlight(
                                    w, category, source="pptx_highlight", confidence=conf,
                                    color_hex=f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"
                                ))
                    escaped = t.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                    if rgb and not _is_neutral_color(*rgb):
                        cat, _ = color_to_category(rgb)
                        if cat and cat in CATEGORY_CSS:
                            html_parts.append(f'<span style="background-color:{CATEGORY_CSS[cat]};">{escaped}</span>')
                        else:
                            html_parts.append(escaped)
                    else:
                        html_parts.append(escaped)
                if line:
                    text_parts.append(line)
                    html_parts.append("<br>")
    html_parts.append("</div>")
    text = "\n".join(text_parts)
    html = "".join(html_parts) if any("span" in p for p in html_parts) else f"<pre>{text}</pre>"
    return text, html, highlights

def _image_bbox_background_color(pil_image, left, top, width, height, margin=3):
    """
    Muestrea el COLOR DE FONDO (background/highlight) alrededor del bbox de una palabra.
    Toma píxeles de un borde exterior (margin px) al bbox y excluye los muy oscuros
    (que serían tinta del texto) y los completamente blancos.
    Retorna (r,g,b) del fondo o None si es neutro/blanco.
    """
    img_w, img_h = pil_image.size
    l = max(0, left - margin)
    t = max(0, top - margin)
    r = min(img_w, left + width + margin)
    b = min(img_h, top + height + margin)
    if r <= l or b <= t:
        return None
    try:
        crop = pil_image.crop((l, t, r, b))
        if crop.mode != "RGB":
            crop = crop.convert("RGB")
        pixels = list(crop.getdata())
        if not pixels:
            return None
        # Excluir tinta (muy oscuros < 80 brightness) y blancos puros (> 245 cada canal)
        bg = [p for p in pixels
              if (p[0] * 0.299 + p[1] * 0.587 + p[2] * 0.114) > 80
              and not (p[0] > 245 and p[1] > 245 and p[2] > 245)]
        if len(bg) < max(4, len(pixels) * 0.1):
            return None
        n = len(bg)
        r_avg = sum(p[0] for p in bg) // n
        g_avg = sum(p[1] for p in bg) // n
        b_avg = sum(p[2] for p in bg) // n
        if _is_neutral_color(r_avg, g_avg, b_avg):
            return None
        return (r_avg, g_avg, b_avg)
    except Exception:
        return None

def extract_text_from_image(file_bytes):
    """Extracts text via OCR and highlights by sampling BACKGROUND color behind each word bbox."""
    image = Image.open(io.BytesIO(file_bytes)).convert("RGB")
    text_plain = pytesseract.image_to_string(image, lang='spa')
    highlights = []
    try:
        data = pytesseract.image_to_data(image, lang='spa', output_type=pytesseract.Output.DICT)
        n = len(data.get("text", []))
        for i in range(n):
            word = (data.get("text") or [])[i]
            if not word or not word.strip():
                continue
            conf = int((data.get("conf") or [0])[i] or 0)
            if conf <= 0:
                continue
            left = int((data.get("left") or [0])[i])
            top = int((data.get("top") or [0])[i])
            width = int((data.get("width") or [0])[i])
            height = int((data.get("height") or [0])[i])
            rgb = _image_bbox_background_color(image, left, top, width, height)
            if rgb:
                category, conf_c = color_to_category(rgb, default_confidence=0.6)
                if category:
                    w_clean = word.strip()
                    if re.match(r'^[\wáéíóúÁÉÍÓÚñÑ]+$', w_clean):
                        highlights.append(normalize_highlight(
                            w_clean, category, source="ocr_background_color", confidence=conf_c,
                            color_hex=f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"
                        ))
    except Exception:
        pass
    html = f"<pre>{text_plain}</pre>"
    return text_plain, html, highlights

def _image_to_highlights_only(pil_image):
    """Dado un PIL Image, devuelve lista de highlights (mismo formato que extract_text_from_image)."""
    try:
        buf = io.BytesIO()
        pil_image.convert("RGB").save(buf, format="PNG")
        _, _, highlights = extract_text_from_image(buf.getvalue())
        return highlights
    except Exception:
        return []

def _merge_highlights_by_word_category(highlight_lists):
    """Fusiona listas de highlights: por (word, expected_category) se queda la mayor confidence."""
    key_to_best = {}
    for lst in highlight_lists:
        for h in lst:
            w = (h.get("word") or "").strip()
            c = h.get("expected_category")
            if not w or not c:
                continue
            key = (w.lower(), c)
            conf = float(h.get("confidence", 0))
            if key not in key_to_best or conf > key_to_best[key].get("confidence", 0):
                key_to_best[key] = dict(h)
    return list(key_to_best.values())

def extract_highlights_from_video(file_bytes, extension, frame_interval_sec=2.0, max_frames=30):
    """
    Extrae frames cada frame_interval_sec segundos, corre OCR+color en cada uno y fusiona highlights.
    """
    merged = []
    video_path = None
    video = None
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=extension) as temp_video:
            temp_video.write(file_bytes)
            video_path = temp_video.name
        video = VideoFileClip(video_path)
        duration = video.duration or 0
        if duration <= 0:
            return merged
        t = 0.0
        frame_highlights = []
        n = 0
        while t < duration and n < max_frames:
            try:
                frame = video.get_frame(t)
                pil_img = Image.fromarray(frame)
                frame_highlights.append(_image_to_highlights_only(pil_img))
                n += 1
            except Exception:
                pass
            t += frame_interval_sec
        merged = _merge_highlights_by_word_category(frame_highlights)
    except Exception:
        pass
    finally:
        if video is not None:
            try:
                video.close()
            except Exception:
                pass
        if video_path and os.path.isfile(video_path):
            try:
                os.remove(video_path)
            except Exception:
                pass
    return merged

def extract_audio_from_video(file_bytes, extension):
    """Extracts audio from a video file and returns the path to the temp audio file."""
    with tempfile.NamedTemporaryFile(delete=False, suffix=extension) as temp_video:
        temp_video.write(file_bytes)
        video_path = temp_video.name
    
    audio_path = video_path.replace(extension, ".mp3")
    try:
        video = VideoFileClip(video_path)
        video.audio.write_audiofile(audio_path, logger=None)
        return audio_path, video_path
    except Exception as e:
        return None, video_path

def process_uploaded_file(uploaded_file):
    """Router for processing different file types. Todos los formatos pueden devolver highlights si hay color detectable."""
    name = uploaded_file.name
    ext = os.path.splitext(name)[1].lower()
    content = uploaded_file.read()

    default_meta = {"titulo_ok": False, "fuente_ok": True, "audio_url": None, "highlight_fallback_reason": None}

    if ext == ".docx":
        text, html, highlights, meta = extract_text_from_docx(content)
        return text, html, None, highlights, meta
    elif ext == ".pdf":
        text, html, highlights = extract_text_from_pdf(content)
        meta = {
            **default_meta, 
            "titulo_ok": detect_titulo_presente(text), 
            "audio_url": extract_audio_url_from_text(text),
            "html_with_highlights": html
        }
        if not highlights:
            meta["highlight_fallback_reason"] = "pdf_sin_sombreado_ni_anotaciones_highlight"
        return text, html, None, highlights, meta
    elif ext in [".pptx", ".ppt"]:
        text, html, highlights = extract_text_from_pptx(content)
        meta = {**default_meta, "titulo_ok": detect_titulo_presente(text), "audio_url": extract_audio_url_from_text(text)}
        if not highlights:
            meta["highlight_fallback_reason"] = "pptx_sin_sombreado_highlight"
        return text, html, None, highlights, meta
    elif ext in [".png", ".jpg", ".jpeg", ".tiff"]:
        text, html, highlights = extract_text_from_image(content)
        meta = {**default_meta, "titulo_ok": detect_titulo_presente(text), "audio_url": extract_audio_url_from_text(text)}
        if not highlights:
            meta["highlight_fallback_reason"] = "imagen_sin_fondo_de_color_detectable"
        return text, html, None, highlights, meta
    elif ext in [".mp4", ".mov", ".avi", ".mkv"]:
        audio_path, video_path = extract_audio_from_video(content, ext)
        highlights = extract_highlights_from_video(content, ext)
        meta = {**default_meta}
        if not highlights:
            meta["highlight_fallback_reason"] = "video_sin_texto_con_color_en_frames"
        return f"[Video File: Audio extracted for transcription at {audio_path}]", None, audio_path, highlights, meta
    elif ext in [".mp3", ".wav", ".m4a"]:
        with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as temp_audio:
            temp_audio.write(content)
        meta = {**default_meta, "highlight_fallback_reason": "formato_audio_sin_color"}
        return f"[Audio File: Saved for transcription at {temp_audio.name}]", None, temp_audio.name, [], meta
    else:
        return None, None, None, [], default_meta
